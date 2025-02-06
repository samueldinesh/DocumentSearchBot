import os
import shutil
import logging
import math
from dotenv import load_dotenv
from fastapi import FastAPI, UploadFile, File, HTTPException, status, Body
from fastapi.middleware.cors import CORSMiddleware

from validate_file import validate_file_size_type  # our updated validation function

# Document processing libraries
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation

# For embedding and vector store
from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain.prompts import ChatPromptTemplate
from langchain_community.vectorstores import FAISS
from langchain.memory import ConversationBufferMemory



# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

load_dotenv()
openai_api_key = os.getenv("OPENAI_API_KEY")
if not openai_api_key:
    logger.error("Error: OPENAI_API_KEY is missing. Please set it in .env file.")
    raise ValueError("OPENAI_API_KEY is missing. Please set it in .env file.")

# Storage paths
DOCUMENT_PATH = "storage"  # ensure this directory exists
FAISS_INDEX_PATH = "faiss_index"

# Initialize OpenAI embeddings
embeddings = OpenAIEmbeddings(api_key=openai_api_key)

# Global variables for vector store and retriever
vectorstore = None
retriever = None

if os.path.exists(FAISS_INDEX_PATH):
    try:
        logger.info("Loading FAISS index...")
        vectorstore = FAISS.load_local(FAISS_INDEX_PATH, embeddings, allow_dangerous_deserialization=True)
        retriever = vectorstore.as_retriever()
    except Exception as e:
        logger.error(f"Failed to load FAISS index: {e}")
else:
    logger.info("FAISS index not found. A new index will be created upon first upload.")

# Initialize LLM and conversation memory
llm = ChatOpenAI(api_key=openai_api_key)

# Define prompt template
prompt_template = ChatPromptTemplate.from_template(
    """
    Context:
    {context}

    Task:
    The user is asking: "{question}"
    - Give the constructive answer from the context and chat history.
    - The answer should be reasoning and logical.
    """
)

# Create FastAPI app
server = FastAPI()

# Add CORS middleware
server.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change "*" to your frontend URL in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


#######################
# Utility Functions
#######################

def clear_vectorstore():
    """
    Clears the in-memory vectorstore and removes any persisted FAISS index files.
    """
    global vectorstore, retriever
    vectorstore = None
    retriever = None
    if os.path.exists(FAISS_INDEX_PATH):
        try:
            # Remove all files in the FAISS_INDEX_PATH directory
            files = os.listdir(FAISS_INDEX_PATH)
            for file in files:
                filepath = os.path.join(FAISS_INDEX_PATH, file)
                if os.path.isfile(filepath):
                    os.remove(filepath)
            logger.info("Cleared FAISS index files.")
        except Exception as e:
            logger.error(f"Error clearing FAISS index: {e}")

def extract_text_from_pdf(file_path: str) -> str:
    text = ""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error processing PDF")
    return text

def extract_text_from_docx(file_path: str) -> str:
    text = ""
    try:
        doc = docx.Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error processing DOCX")
    return text

def extract_text_from_excel(file_path: str) -> str:
    text = ""
    try:
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text += df.to_string() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from Excel: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error processing Excel")
    return text

def extract_text_from_pptx(file_path: str) -> str:
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error processing PPTX")
    return text

def chunk_text(text: str, chunk_size: int = 3500) -> list:
    """
    Splits text into chunks based on word count.
    (For more accurate token counts, a tokenizer can be integrated later.)
    """
    words = text.split()
    chunks = []
    total_words = len(words)
    num_chunks = math.ceil(total_words / chunk_size)
    for i in range(num_chunks):
        start = i * chunk_size
        end = start + chunk_size
        chunk = " ".join(words[start:end])
        chunks.append(chunk)
    logger.info(f"Text split into {len(chunks)} chunks.")
    return chunks

def update_faiss_index(chunks: list, metadata: dict = None):
    global vectorstore, retriever
    try:
        if vectorstore is None:
            vectorstore = FAISS.from_texts(chunks, embeddings, metadatas=[metadata] * len(chunks))
            logger.info("Created new FAISS index.")
        else:
            vectorstore.add_texts(chunks, metadatas=[metadata] * len(chunks))
            logger.info("Updated existing FAISS index with new document chunks.")
        vectorstore.save_local(FAISS_INDEX_PATH)
        # Update retriever reference
        retriever = vectorstore.as_retriever()
    except Exception as e:
        logger.error(f"Error updating FAISS index: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error updating vector store")

def process_document(file_path: str, filename: str) -> None:
    """
    Clears the current vectorstore, then processes the new document by extracting text,
    chunking it, and storing its embeddings in a fresh FAISS index.
    """
    # Clear any existing vectorstore data
    clear_vectorstore()

    ext = os.path.splitext(filename)[1].lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext in [".doc", ".docx"]:
        text = extract_text_from_docx(file_path)
    elif ext in [".xls", ".xlsx"]:
        text = extract_text_from_excel(file_path)
    elif ext in [".ppt", ".pptx"]:
        text = extract_text_from_pptx(file_path)
    else:
        logger.error("Unsupported file extension for processing")
        raise HTTPException(status_code=status.HTTP_415_UNSUPPORTED_MEDIA_TYPE, detail="Unsupported file extension for processing")
    
    if not text.strip():
        logger.error("No text extracted from the document.")
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="No text extracted from the document")
    
    chunks = chunk_text(text)
    metadata = {"filename": filename}
    update_faiss_index(chunks, metadata)
    logger.info(f"Processed document {filename} and updated FAISS index with {len(chunks)} chunks.")

##############################
# Endpoints
##############################

@server.post("/uploadfile/")
async def create_upload_file(file: UploadFile = File(...)):
    # Validate file size and type:
    try:
        file.file.seek(0)
        validate_file_size_type(file)
        file.file.seek(0)
    except Exception as e:
        logger.error(f"File validation error: {e}")
        raise

    # Save file to disk
    try:
        os.makedirs(DOCUMENT_PATH, exist_ok=True)
        file_location = os.path.join(DOCUMENT_PATH, file.filename)
        with open(file_location, "wb+") as file_object:
            shutil.copyfileobj(file.file, file_object)
        logger.info(f"File {file.filename} uploaded successfully to {file_location}.")
    except Exception as e:
        logger.error(f"Error saving file: {e}")
        raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="File upload failed")
    
    # Process the document (this clears the vectorstore first)
    try:
        process_document(file_location, file.filename)
    except Exception as e:
        logger.error(f"Error processing document {file.filename}: {e}")
        raise

    return {"message": f"File {file.filename} uploaded and processed successfully."}

@server.get("/getfile")
async def getFile():
    if os.path.exists(DOCUMENT_PATH):
        files = [f for f in os.listdir(DOCUMENT_PATH) if os.path.isfile(os.path.join(DOCUMENT_PATH, f))]
        logger.info(f"Files found: {files}")
        return {"files": files}
    return {"files": "Not found"}

@server.delete("/deletefile/{filename}")
async def delete_file_and_data(filename: str):
    """
    Delete the uploaded file from disk and remove its chunks from the FAISS index.
    """
    # Delete the file from disk
    file_location = os.path.join(DOCUMENT_PATH, filename)
    if os.path.exists(file_location):
        try:
            os.remove(file_location)
            logger.info(f"File {filename} deleted from disk.")
        except Exception as e:
            logger.error(f"Error deleting file {filename}: {e}")
            raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error deleting file")
    else:
        logger.error(f"File {filename} not found on disk.")
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File not found")

    # Remove corresponding chunks from the vector store
    if vectorstore is None:
        logger.info("No FAISS index exists; nothing to delete in vector store.")
    else:
        try:
            logger.info("Clearing FAISS index as part of deletion process.")
            clear_vectorstore()
            logger.info(f"Cleared FAISS index after deleting file {filename}.")
        except Exception as e:
            logger.error(f"Error updating vector store after deletion: {e}")
            raise HTTPException(status_code=status.HTTP_500_INTERNAL_SERVER_ERROR, detail="Error updating vector store")
    
    return {"message": f"File {filename} and its vector store data deleted successfully."}

@server.post("/chat")
async def chat_with_bot(user_message: str = Body(..., embed=True)):
    logger.info(f"Received user_message: {user_message}")
    # Retrieve relevant documents
    if retriever is None:
        context = "No context available."
    else:
        retrieved_docs = retriever.invoke(user_message)
        context = retrieved_docs[:3] if retrieved_docs else "No direct reference found."
   
    # Generate AI response
    prompt = prompt_template.format(
        context=context,
        question=user_message
    )
    response = llm.invoke(prompt)
    
    return {"response": response.content}
