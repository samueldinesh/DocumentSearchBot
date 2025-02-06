import os
import math
import logging
from typing import List
from fastapi import HTTPException
from app.config import DOCUMENT_PATH
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation

logger = logging.getLogger(__name__)

class DocumentProcessor:
    @staticmethod
    def extract_text(file_path: str, filename: str) -> str:
        ext = os.path.splitext(filename)[1].lower()
        try:
            if ext == ".pdf":
                return DocumentProcessor._extract_pdf(file_path)
            elif ext in [".doc", ".docx"]:
                return DocumentProcessor._extract_docx(file_path)
            elif ext in [".xls", ".xlsx"]:
                return DocumentProcessor._extract_excel(file_path)
            elif ext in [".ppt", ".pptx"]:
                return DocumentProcessor._extract_pptx(file_path)
            raise HTTPException(status_code=415, detail="Unsupported file type")
        except Exception as e:
            logger.error(f"Error processing {filename}: {str(e)}")
            raise HTTPException(status_code=500, detail=f"Error processing {filename}")

    @staticmethod
    def _extract_pdf(file_path: str) -> str:
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += page.extract_text() or ""
        return text

    @staticmethod
    def _extract_docx(file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join(para.text for para in doc.paragraphs)

    @staticmethod
    def _extract_excel(file_path: str) -> str:
        text = []
        xl = pd.ExcelFile(file_path)
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name)
            text.append(df.to_string())
        return "\n".join(text)

    @staticmethod
    def _extract_pptx(file_path: str) -> str:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 3500) -> List[str]:
        words = text.split()
        return [" ".join(words[i:i+chunk_size]) for i in range(0, len(words), chunk_size)]